"""Extractor PPTX: diapositiva a diapositiva con elementos ordenados visualmente."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from mmi.index.content_hash import slide_content_hash
from mmi.ingest.pptx_models import (
    PresentationExtract,
    SlideElement,
    SlideQuality,
    SlideRecord,
)
from mmi.ingest.pptx_visual import maybe_describe_visual


def _shape_position(shape) -> tuple[int, int]:
    return (int(shape.top or 0), int(shape.left or 0))


def _table_to_markdown(headers: list[str], rows: list[list[str]]) -> str:
    if not headers and not rows:
        return ""
    hdr = headers or (rows[0] if rows else [])
    body = rows if headers else rows[1:]
    if not hdr:
        return ""
    lines = [
        "| " + " | ".join(hdr) + " |",
        "| " + " | ".join("---" for _ in hdr) + " |",
    ]
    for row in body:
        padded = row + [""] * (len(hdr) - len(row))
        lines.append("| " + " | ".join(padded[: len(hdr)]) + " |")
    return "\n".join(lines)


def _extract_table(shape) -> SlideElement | None:
    try:
        table = shape.table
    except (AttributeError, ValueError):
        return None
    rows_raw: list[list[str]] = []
    for row in table.rows:
        rows_raw.append([(cell.text or "").strip() for cell in row.cells])
    if not rows_raw or not any(any(c for c in r) for r in rows_raw):
        return None
    headers = rows_raw[0]
    body = rows_raw[1:] if len(rows_raw) > 1 else []
    md = _table_to_markdown(headers, body if body else rows_raw)
    return SlideElement(
        kind="table",
        order=0,
        headers=headers,
        rows=body if body else rows_raw,
        markdown=md,
        text=md,
    )


def _extract_chart(shape, order: int) -> SlideElement | None:
    if not getattr(shape, "has_chart", False):
        return None
    try:
        chart = shape.chart
    except (AttributeError, ValueError):
        return None
    title = ""
    try:
        if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
    except (AttributeError, ValueError):
        pass

    axes: dict[str, str] = {}
    try:
        cat_axis = chart.category_axis
        if cat_axis and cat_axis.has_title:
            axes["x"] = cat_axis.axis_title.text_frame.text.strip()
    except (AttributeError, ValueError):
        pass
    try:
        val_axis = chart.value_axis
        if val_axis and val_axis.has_title:
            axes["y"] = val_axis.axis_title.text_frame.text.strip()
    except (AttributeError, ValueError):
        pass

    legend: list[str] = []
    series: list[dict[str, Any]] = []
    try:
        for ser in chart.series:
            name = getattr(ser, "name", "") or ""
            if name:
                legend.append(str(name))
            values: list[Any] = []
            try:
                values = list(ser.values) if ser.values is not None else []
            except (AttributeError, ValueError, TypeError):
                pass
            series.append({"name": str(name), "values": values})
    except (AttributeError, ValueError):
        pass

    parts = [f"Gráfico: {title}" if title else "Gráfico"]
    if axes:
        parts.append(f"Ejes: {axes}")
    if legend:
        parts.append(f"Leyenda: {', '.join(legend)}")
    for s in series:
        vals = s.get("values") or []
        if vals:
            parts.append(f"{s.get('name', 'Serie')}: {vals}")

    return SlideElement(
        kind="chart",
        order=order,
        chart_title=title or None,
        axes=axes,
        legend=legend,
        series=series,
        text="\n".join(parts),
    )


def _extract_text_shape(shape, order: int, kind: str = "text_box") -> SlideElement | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    text = (shape.text or "").strip()
    if not text:
        return None
    return SlideElement(kind=kind, order=order, text=text)


def _extract_image(shape, order: int, slide_number: int) -> SlideElement:
    media_ref = f"slide_{slide_number}_img_{order}"
    return SlideElement(
        kind="image",
        order=order,
        media_ref=media_ref,
        needs_visual_analysis=True,
    )


def _slide_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except (AttributeError, ValueError):
        pass
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            if shape.is_placeholder and shape.placeholder_format.type in (1, 3):  # TITLE, CENTER_TITLE
                t = (shape.text or "").strip()
                if t:
                    return t
        except (AttributeError, ValueError):
            continue
    return ""


def _speaker_notes(slide) -> str:
    try:
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            return (notes_slide.notes_text_frame.text or "").strip()
    except (AttributeError, ValueError):
        pass
    return ""


def _master_context(slide) -> str:
    parts: list[str] = []
    try:
        for ph in slide.placeholders:
            try:
                ptype = ph.placeholder_format.type
                if ptype in (2, 4, 5):  # BODY, SUBTITLE, DATE
                    continue
                if ptype in (6, 7, 8, 9, 10):  # FOOTER, HEADER, etc.
                    t = (ph.text or "").strip()
                    if t and t not in parts:
                        parts.append(t)
            except (AttributeError, ValueError):
                continue
    except (AttributeError, ValueError):
        pass
    return " · ".join(parts)


def _extract_slide_elements(slide, slide_number: int) -> list[SlideElement]:
    elements: list[SlideElement] = []
    shapes = sorted(slide.shapes, key=_shape_position)
    order = 0
    for shape in shapes:
        if getattr(shape, "has_table", False):
            el = _extract_table(shape)
            if el:
                el.order = order
                elements.append(el)
                order += 1
            continue
        if getattr(shape, "has_chart", False):
            el = _extract_chart(shape, order)
            if el:
                elements.append(el)
                order += 1
            continue
        if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
            elements.append(_extract_image(shape, order, slide_number))
            order += 1
            continue
        is_title = False
        try:
            is_title = bool(
                shape.is_placeholder
                and shape.placeholder_format.type in (1, 3)
            )
        except (AttributeError, ValueError):
            pass
        el = _extract_text_shape(shape, order, kind="title" if is_title else "text_box")
        if el:
            elements.append(el)
            order += 1
    return elements


def _assess_slide_quality(slide: SlideRecord) -> SlideQuality:
    has_text = any(e.text.strip() or e.markdown.strip() for e in slide.elements)
    has_table = any(e.kind == "table" and e.markdown for e in slide.elements)
    has_chart = any(e.kind == "chart" and e.text for e in slide.elements)
    has_notes = bool(slide.speaker_notes.strip())
    has_visual = bool(slide.visual_summary.strip())

    if has_text or has_table or has_chart or has_notes:
        only_images = all(e.kind == "image" for e in slide.elements) and not has_notes
        if only_images and not has_visual:
            return "review"
        return "pass"
    if slide.elements and all(e.kind == "image" for e in slide.elements):
        return "review"
    return "reject"


def _build_section_map(prs) -> dict[int, str]:
    """Mapea índice de slide (1-based) → título de sección."""
    mapping: dict[int, str] = {}
    try:
        sections = prs.slide_sections
        for section in sections:
            name = (section.name or "").strip()
            for slide in section.slides:
                mapping[slide.slide_id] = name
    except (AttributeError, ValueError):
        pass
    return mapping


class PptxAdapter:
    """Extrae presentación a slides estructurados."""

    def extract(
        self,
        path: Path,
        *,
        file_hash: str | None = None,
        presentation_id: str | None = None,
        version_id: str | None = None,
        document_key: str | None = None,
    ) -> PresentationExtract:
        from pptx import Presentation

        from mmi.index.chunking import file_sha256

        path = Path(path)
        file_hash = file_hash or file_sha256(path)
        prs = Presentation(str(path))
        notes: list[str] = []

        section_by_id = _build_section_map(prs)

        slides: list[SlideRecord] = []
        current_section = ""

        for idx, slide in enumerate(prs.slides, start=1):
            sid = slide.slide_id
            section_title = section_by_id.get(sid, current_section)
            if section_by_id.get(sid):
                current_section = section_by_id[sid]

            title = _slide_title(slide)
            elements = _extract_slide_elements(slide, idx)
            if not title or title == f"Diapositiva {idx}":
                for el in elements:
                    if el.kind in {"text_box", "title"} and el.text:
                        first_line = el.text.split("\n")[0].strip()
                        if first_line and len(first_line) <= 100:
                            title = first_line
                            break
            if not title:
                title = f"Diapositiva {idx}"
            speaker = _speaker_notes(slide)
            master = _master_context(slide)

            visual_parts: list[str] = []
            for el in elements:
                if el.needs_visual_analysis:
                    desc = maybe_describe_visual(el, path, idx)
                    if desc:
                        el.description = desc
                        visual_parts.append(desc)
                elif el.kind == "chart" and el.text:
                    visual_parts.append(el.text)

            record = SlideRecord(
                presentation_id=presentation_id,
                version_id=version_id,
                slide_number=idx,
                slide_title=title or f"Diapositiva {idx}",
                section_title=section_title or current_section,
                elements=elements,
                speaker_notes=speaker,
                visual_summary="; ".join(visual_parts),
                source_location={
                    "file": path.name,
                    "slide": idx,
                    "document_key": document_key or "",
                    "master_context": master,
                },
            )
            record.slide_content_hash = slide_content_hash(record)
            record.extraction_quality = _assess_slide_quality(record)
            slides.append(record)

        pass_count = sum(1 for s in slides if s.extraction_quality == "pass")
        review_count = sum(1 for s in slides if s.extraction_quality == "review")
        total = len(slides)

        if total == 0:
            quality = "reject"
            notes.append("Presentación sin diapositivas")
        elif pass_count == 0:
            quality = "reject" if review_count == 0 else "review"
            notes.append(f"0/{total} diapositivas con contenido indexable")
        elif review_count > 0:
            quality = "pass" if pass_count >= total * 0.5 else "review"
            notes.append(f"{review_count}/{total} diapositivas solo visuales (revisar)")
        else:
            quality = "pass"

        return PresentationExtract(
            source_path=str(path.resolve()),
            file_hash=file_hash,
            slides=slides,
            quality=quality,
            notes=notes,
            meta={
                "format": "pptx",
                "slide_count": total,
                "slides_pass": pass_count,
                "slides_review": review_count,
                "slides_reject": sum(1 for s in slides if s.extraction_quality == "reject"),
            },
        )

    def to_extracted_document(self, extract: PresentationExtract):
        """Convierte a ExtractedDocument para compatibilidad con ports."""
        from mmi.ingest.ports import ExtractedDocument

        parts = [f"# {Path(extract.source_path).name}", ""]
        for slide in extract.slides:
            parts.append(f"## Diapositiva {slide.slide_number}: {slide.slide_title}")
            if slide.section_title:
                parts.append(f"_Sección: {slide.section_title}_")
            parts.append("")
            for el in slide.elements:
                if el.markdown:
                    parts.append(el.markdown)
                elif el.text:
                    parts.append(el.text)
            if slide.speaker_notes:
                parts.append(f"_Notas: {slide.speaker_notes}_")
            parts.append("")

        return ExtractedDocument(
            markdown="\n".join(parts),
            quality=extract.quality,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            source_path=extract.source_path,
            notes=extract.notes,
            meta=extract.meta,
        )


def save_slides_json(extract: PresentationExtract, target_dir: Path) -> Path:
    target_dir.mkdir(parents=True, exist_ok=True)
    out = target_dir / "slides.json"
    out.write_text(
        json.dumps(extract.slides_to_json(), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return out


def load_slides_json(path: Path) -> list[SlideRecord]:
    data = json.loads(path.read_text(encoding="utf-8"))
    return [SlideRecord.from_dict(item) for item in data]


def load_or_extract(
    path: Path,
    extract_dir: Path | None = None,
    **kwargs,
) -> PresentationExtract:
    """Usa slides.json en cache si file_hash coincide."""
    from mmi.index.chunking import file_sha256

    path = Path(path)
    file_hash = file_sha256(path)
    if extract_dir:
        slides_path = extract_dir / "slides.json"
        meta_path = extract_dir / "extracted.json"
        if slides_path.exists() and meta_path.exists():
            try:
                meta = json.loads(meta_path.read_text(encoding="utf-8"))
                if meta.get("file_hash") == file_hash:
                    slides = load_slides_json(slides_path)
                    return PresentationExtract(
                        source_path=str(path.resolve()),
                        file_hash=file_hash,
                        slides=slides,
                        quality=meta.get("quality", "review"),
                        notes=list(meta.get("notes") or []),
                        meta=dict(meta.get("meta") or {}),
                    )
            except (OSError, json.JSONDecodeError, KeyError):
                pass
    adapter = PptxAdapter()
    return adapter.extract(path, file_hash=file_hash, **kwargs)
