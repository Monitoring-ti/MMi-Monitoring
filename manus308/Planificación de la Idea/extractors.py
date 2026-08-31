#!/usr/bin/env python3
"""
MMI — Fase 1 · Extractores de contenido por formato.

Cada extractor devuelve una lista de "bloques" de texto con metadatos de
procedencia (página / hoja / slide) que luego el chunker segmenta.

  - PdfExtractor    : texto nativo con pdfplumber (páginas con capa de texto).
                      Si una página no tiene texto, marca needs_ocr=True.
  - XlsxExtractor   : extracción tabular con openpyxl preservando celdas
                      vacías como null y cabeceras repetidas por fila.
  - PptxExtractor   : texto slide a slide con python-pptx, incluyendo notas
                      del orador (speaker notes).
"""
from __future__ import annotations

import os
from dataclasses import dataclass, field
from typing import Optional


@dataclass
class Block:
    """Bloque de texto extraído con su procedencia."""
    text: str
    page: Optional[int] = None          # PDF
    sheet: Optional[str] = None         # XLSX
    row: Optional[int] = None           # XLSX
    slide: Optional[int] = None         # PPTX
    notes: Optional[str] = None         # PPTX speaker notes
    needs_ocr: bool = False             # PDF escaneado sin capa de texto
    meta: dict = field(default_factory=dict)


# ----------------------------------------------------------------------------
# PDF
# ----------------------------------------------------------------------------

class PdfExtractor:
    """Extrae texto nativo página a página. Detecta páginas escaneadas."""

    def extract(self, path: str) -> list[Block]:
        import pdfplumber
        blocks: list[Block] = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text:
                    blocks.append(Block(text=text, page=i))
                else:
                    # Página sin capa de texto -> candidata a OCR
                    blocks.append(Block(text="", page=i, needs_ocr=True))
        return blocks


# ----------------------------------------------------------------------------
# XLSX (tabular, null preservados)
# ----------------------------------------------------------------------------

class XlsxExtractor:
    """Convierte cada fila en una línea 'columna: valor' preservando nulls.

    Repite la cabecera de cada hoja para dar contexto a cada fila (mejor
    retrieval). Las celdas vacías se representan como 'null' explícito.
    """

    def extract(self, path: str) -> list[Block]:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        blocks: list[Block] = []
        for ws in wb.worksheets:
            rows = [[self._cell(c) for c in r] for r in ws.iter_rows(values_only=True)]
            header_idx, data_start = self._locate_table(rows)
            if header_idx is None:
                continue
            headers = self._build_headers(rows, header_idx)
            # Omitir columnas vacías en TODAS las filas de datos (margen
            # izquierdo típico de formularios).
            used = [False] * len(headers)
            for r_idx in range(data_start, len(rows)):
                for c, v in enumerate(rows[r_idx]):
                    if c < len(headers) and v:
                        used[c] = True
            keep = [c for c in range(len(headers)) if used[c] or headers[c]]
            last_vals = [""] * len(headers)
            for r_idx in range(data_start, len(rows)):
                raw = rows[r_idx]
                if not any(raw):
                    continue  # fila vacía -> se omite
                parts = []
                for c in keep:
                    h = headers[c]
                    val = raw[c] if c < len(raw) else ""
                    if val:
                        last_vals[c] = val
                        shown = val
                    else:
                        shown = last_vals[c] if last_vals[c] else "null"
                    label = h if h else f"col_{c + 1}"
                    parts.append(f"{label}: {shown}")
                line = " | ".join(parts).strip()
                if line:
                    blocks.append(Block(text=line, sheet=ws.title, row=r_idx + 1,
                                        meta={"headers": headers}))
        wb.close()
        return blocks

    # -- utilidades ---------------------------------------------------------

    @staticmethod
    def _filled(row) -> int:
        return sum(1 for c in row if c)

    def _locate_table(self, rows):
        """Localiza la tabla. Devuelve (header_idx, data_start).

        La cabecera es la fila con más columnas de texto. Si la fila siguiente
        parece una SUB-cabecera (p.ej. 'SI | NO | N/A' bajo 'CUMPLIMIENTO'),
        los datos empiezan una fila más abajo y las cabeceras se combinan."""
        n = len(rows)
        best_header, best_score = None, 0
        for i in range(n):
            score = self._filled(rows[i])
            if score > best_score:
                best_header, best_score = i, score
        if best_header is None or best_score < 2:
            return None, None
        # Si la fila siguiente es una sub-cabecera (pocas celdas cortas tipo
        # SI/NO/N/A), los datos empiezan una fila después.
        nxt = best_header + 1
        if nxt < n and self._is_subheader(rows[nxt]):
            return best_header, best_header + 2
        return best_header, best_header + 1

    @staticmethod
    def _is_subheader(row) -> bool:
        """Detecta una fila de sub-cabecera tipo 'SI | NO | N/A' (valores muy
        cortos y categóricos)."""
        vals = [c for c in row if c]
        if not vals:
            return False
        short = sum(1 for v in vals if len(v) <= 4)
        return len(vals) >= 2 and short / len(vals) >= 0.6

    def _build_headers(self, rows, header_idx: int) -> list[str]:
        """Construye las cabeceras combinando la fila de grupo (si existe) con
        la fila de cabecera, y la sub-cabecera si la hay. Forward-fill horizontal
        para celdas combinadas. Deduplica cabeceras repetidas con sufijo."""
        def ff(row):
            out, last = [], ""
            for c in row:
                if c:
                    last = c
                out.append(last)
            return out

        main = ff(rows[header_idx])
        # Nivel de grupo (fila anterior) solo si es densa
        if header_idx > 0 and self._filled(rows[header_idx - 1]) >= 2:
            top = ff(rows[header_idx - 1])
            main = [f"{t} > {m}".strip(" >") if t and t != m else (m or t)
                    for t, m in zip(top, main)]
        # Sub-cabecera (fila siguiente) si existe
        nxt = header_idx + 1
        if nxt < len(rows) and self._is_subheader(rows[nxt]):
            sub = rows[nxt]
            main = [f"{m} > {s}".strip(" >") if s else m
                    for m, s in zip(main, sub + [""] * (len(main) - len(sub)))]
        # Asignar nombres semánticos a columnas sin cabecera propia según su
        # posición y contenido típico (N°, Aspecto, etc.). Luego deduplicar.
        out = []
        seen: dict[str, int] = {}
        positional = 0
        for h in main:
            if not h:
                out.append("")
                continue
            if h in seen:
                seen[h] += 1
                # En vez de 'X_2', dar un nombre por rol si es la 2ª columna
                # (típicamente el texto del aspecto/pregunta).
                out.append(f"{h}_{seen[h]}")
            else:
                seen[h] = 1
                out.append(h)
        return out

    @staticmethod
    def _cell(v) -> str:
        if v is None:
            return ""
        return str(v).strip()


# ----------------------------------------------------------------------------
# PPTX (slide a slide + speaker notes)
# ----------------------------------------------------------------------------

class PptxExtractor:
    """Extrae el texto de cada slide y sus notas del orador."""

    def extract(self, path: str) -> list[Block]:
        from pptx import Presentation
        prs = Presentation(path)
        blocks: list[Block] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        texts.append(" | ".join(cells))
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            body = "\n".join(texts).strip()
            if body or notes:
                blocks.append(Block(text=body, slide=i, notes=notes or None))
        return blocks


# ----------------------------------------------------------------------------
# DOCX (párrafos + tablas)
# ----------------------------------------------------------------------------

class DocxExtractor:
    """Extrae párrafos y tablas de un DOCX. Cada párrafo con texto es un
    bloque; cada fila de tabla se serializa como 'col: val | col: val'."""

    def extract(self, path: str) -> list[Block]:
        import docx
        doc = docx.Document(path)
        blocks: list[Block] = []
        # Párrafos
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                blocks.append(Block(text=t))
        # Tablas
        for ti, table in enumerate(doc.tables, start=1):
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    blocks.append(Block(text=line, sheet=f"tabla_{ti}"))
        return blocks


# ----------------------------------------------------------------------------
# Dispatcher por extensión
# ----------------------------------------------------------------------------

_EXTRACTORS = {
    ".pdf": PdfExtractor,
    ".xlsx": XlsxExtractor,
    ".pptx": PptxExtractor,
    ".docx": DocxExtractor,
}


def extract(path: str) -> list[Block]:
    ext = os.path.splitext(path)[1].lower()
    cls = _EXTRACTORS.get(ext)
    if cls is None:
        raise ValueError(f"Formato no soportado: {ext} ({path})")
    return cls().extract(path)


if __name__ == "__main__":
    import sys
    p = sys.argv[1]
    bl = extract(p)
    total_chars = sum(len(b.text) for b in bl)
    ocr = sum(1 for b in bl if b.needs_ocr)
    print(f"{os.path.basename(p)}: {len(bl)} bloques, {total_chars} chars, "
          f"{ocr} páginas necesitan OCR")
